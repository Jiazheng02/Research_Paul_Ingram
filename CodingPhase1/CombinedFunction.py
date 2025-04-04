import re
import pandas as pd
from pptx import Presentation
from NodesLinesMatching import matching_lines_nodes

def process_pptx_to_dataframe(pptx_path, id_path):
    """
    Extracts person names and matched edges (from → to nodes) for each slide in a PPTX file,
    ensuring all slides are included even if no edges exist.

    Args:
        pptx_path (str): Path to the PowerPoint (.pptx) file.
        id_path (str): Path to the Excel file mapping names to id_ego.

    Returns:
        pd.DataFrame: DataFrame containing all slides, including those without extracted edges.
    """
    prs = Presentation(pptx_path)
    data = []  # Store slide information and edges

    for slide_num, slide in enumerate(prs.slides, start=1):
        sorted_edges, respondent, node_count, line_count = matching_lines_nodes(slide, slide_num)
        void_edge = int(line_count == 0)  # Indicator for missing edges

        if not respondent:
            continue
            
        # Store respondent and slide info (even if no edges exist)
        if node_count == 0 or line_count == 0:
            data.append([respondent, None, None, void_edge, slide_num])  # No nodes or edges, placeholders for "From" & "To"
        else:
            for from_text, to_text in sorted_edges:
                data.append([respondent, from_text, to_text, void_edge, slide_num])

    # Create DataFrame
    df = pd.DataFrame(data, columns=['Person Name', 'From', 'To', 'NoMeaningfulEdges', 'Slide Number'])

    # Load ID mapping and apply
    id_df_map = pd.read_excel(id_path).set_index('full name')['id_ego'].to_dict()
    df.insert(0, 'id_ego', df['Person Name'].map(id_df_map))  # Insert ID at the first column

    def clean_dataframe(df):
        """Remove illegal characters from all string values in the DataFrame."""
        if not isinstance(df, pd.DataFrame):  # Ensure df is a DataFrame
            raise TypeError("Expected a pandas DataFrame but got something else")
        
        # Convert the column to string type to ensure consistent processing of all values.
        # Use regex to remove any illegal ASCII control characters (range \x00-\x1F and \x7F)
        # that can cause issues when saving to Excel.
        # The condition `if pd.notna(x) else x` ensures NaN values remain unchanged.
        for col in df.select_dtypes(include=['object']).columns:  # Select only string columns
            df[col] = df[col].astype(str).apply(lambda x: re.sub(r'[\x00-\x1F\x7F]', ' ', x) if pd.notna(x) else x)
            
        return df
    
    df = clean_dataframe(df)
    
    return df