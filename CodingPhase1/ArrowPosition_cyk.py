from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.connector import Connector
import matplotlib.pyplot as plt
import math


def get_arrowhead_info(shape):
    """
    Retrieve arrowhead types from the shape's XML.
    """
    xml_elem = shape.element
    ns = xml_elem.nsmap
    ln = xml_elem.find(".//a:ln", ns)
    head_type = None
    tail_type = None
    if ln is not None:
        head = ln.find("a:headEnd", ns)
        tail = ln.find("a:tailEnd", ns)
        if head is not None:
            head_type = head.get("type")
        if tail is not None:
            tail_type = tail.get("type")
    return head_type, tail_type


def rotate_line(begin_x, begin_y, end_x, end_y, rot_degrees):
    mid_x = (begin_x + end_x) / 2
    mid_y = (begin_y + end_y) / 2

    theta = math.radians(rot_degrees)

    def rotate_point(x, y):
        translated_x = x - mid_x
        translated_y = y - mid_y

        rotated_x = translated_x * math.cos(theta) - translated_y * math.sin(theta)
        rotated_y = translated_x * math.sin(theta) + translated_y * math.cos(theta)

        return rotated_x + mid_x, rotated_y + mid_y

    new_begin_x, new_begin_y = rotate_point(begin_x, begin_y)
    new_end_x, new_end_y = rotate_point(end_x, end_y)

    return new_begin_x, new_begin_y, new_end_x, new_end_y


def main():
    prs = Presentation("test.pptx")
    plt.gca().invert_yaxis()
    for i, slide in enumerate(prs.slides, start=1):
        print(f"Slide {i}:")
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                shape: Connector = shape
                begin_x = shape.begin_x
                begin_y = shape.begin_y
                end_x = shape.end_x
                end_y = shape.end_y
                rot = shape.element.rot
                begin_x, begin_y, end_x, end_y = rotate_line(
                    begin_x, begin_y, end_x, end_y, rot
                )
                print(f"  Line from ({begin_x}, {begin_y}) to ({end_x}, {end_y})")
                arrow_head, arrow_tail = get_arrowhead_info(shape)
                print(f"    Arrowhead: {arrow_head}, Arrowtail: {arrow_tail}")
                plt.plot([begin_x, end_x], [begin_y, end_y], "k-")
        break
    plt.show()


if __name__ == "__main__":
    main()
