from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.connector import Connector

def extract_lines(slide):
    """
    Counts and extracts the positions of all line shapes (including arrows and connectors)
    from a given PowerPoint slide.

    Args:
        slide (Slide): A PowerPoint slide object from a .pptx file.

    Returns:
        tuple:
            - line_count (int): The total number of lines found on the slide.
            - lines (dict): A dictionary where:
                - Key (int): The index of the line shape within the slide.
                - Value (tuple): The start and end coordinates of the line in the format (x1, y1, x2, y2), where:
                    - (x1, y1): The starting point of the line.
                    - (x2, y2): The ending point of the line.
    """
    line_count = 0  # Counter for the number of detected lines
    lines = {}  # Dictionary to store extracted line positions

    # Iterate through all shapes in the slide
    for i, shape in enumerate(slide.shapes):
        # Check if the shape is a line (includes arrows and connectors)
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            line_count += 1  # Increment line count

            # Explicitly cast shape to Connector type for clarity
            shape: Connector = shape  

            # Extract start and end coordinates of the line
            x1, y1 = shape.begin_x, shape.begin_y  # Start point
            x2, y2 = shape.end_x, shape.end_y  # End point

            # Store the extracted line coordinates in the dictionary
            lines[i] = (x1, y1, x2, y2)

    return lines, line_count  # Return the total count and dictionary of lines