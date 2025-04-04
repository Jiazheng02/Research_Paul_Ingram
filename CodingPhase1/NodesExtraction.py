from pptx import Presentation

def extract_node_boxes(slide, slide_num):
    """
    Extract all text boxes from a single slide and sort by position.
    Also identifies the Name (ID) of the slide as the topmost-leftmost text.
    
    Args:
        slide: The slide object from pptx.
        slide_num (int): The slide number.

    Returns:
        tuple: (List of node_boxes, Name of the slide)
    """
    node_boxes = []
    respondent = None # Store the name
    
    def extract_text_from_shape(shape, slide_num):
        """Recursively extract text from shapes, including grouped shapes."""
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                top, left, right, bottom = shape.top, shape.left, shape.left + shape.width, shape.top + shape.height
                node_boxes.append((slide_num, top, left, right, bottom, text))

        # Handle GroupShape elements (shapes that contain other shapes)
        if shape.shape_type == 6:  # MSO_SHAPE.GROUP
            for sub_shape in shape.shapes:
                extract_text_from_shape(sub_shape, slide_num)
    
    # Extract text from all shapes on the slide
    for shape in slide.shapes:
        extract_text_from_shape(shape, slide_num)

    # Identify the name (ID) → The topmost-leftmost text box
    if node_boxes:
        name_box = min(node_boxes, key=lambda x: (x[1], x[2]))  # Sort by top, then left by (slide_num, shape.top, shape.left, text)
        if not name_box:  # Ensure name_box is not None
            return None
        respondent = name_box[5]  # Extract the text as the slide name
        node_boxes.remove(name_box)  # Remove name from nodes to avoid duplication
    
    # Count valid text nodes (after removing respondent)
    node_count = len(node_boxes)
    
    # Sort nodes by (slide number, vertical position, horizontal position)
    node_boxes.sort(key=lambda x: (x[0], x[1], x[2]))

    return node_boxes, respondent, node_count