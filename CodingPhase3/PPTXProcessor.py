import io
import re
import cv2
import pandas as pd
import numpy as np
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.connector import Connector
from paddleocr import PaddleOCR
from ImageProcessor import ImageProcessor

class PPTXProcessor:
    """
    Extracts text boxes, lines, and relationships from a PowerPoint (.pptx) file.
    Uses PPT text if possible; falls back to ImageProcessor when no lines are detected.
    """

    def __init__(self, pptx_path, id_path):
        """
        Initializes the extractor.

        :param pptx_path: Path to the PowerPoint file (.pptx).
        :param id_path: Path to the Excel file mapping names to 'id_ego'.
        :param image_processor: ImageProcessor instance for OCR-based extraction.
        """
        self.pptx_path = pptx_path
        self.id_path = id_path
        self.prs = Presentation(pptx_path)

    def extract_node_boxes(self, slide, slide_num):
        """
        Extracts text boxes from a PPT slide and identifies the respondent (title).

        :param slide: The slide object.
        :param slide_num: The slide number.
        :return: Tuple (List of node_boxes, Respondent Name, Node Count)
        """
        node_boxes = []
        respondent = None  

        def extract_text_from_shape(shape):
            """Recursively extracts text from shapes, including grouped shapes."""
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    top, left, right, bottom = shape.top, shape.left, shape.left + shape.width, shape.top + shape.height
                    node_boxes.append((slide_num, top, left, right, bottom, text))

            # Handle GroupShapes containing multiple text boxes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    extract_text_from_shape(sub_shape)

        for shape in slide.shapes:
            extract_text_from_shape(shape)

        # Identify respondent (top-leftmost text box)
        if node_boxes:
            name_box = min(node_boxes, key=lambda x: (x[1], x[2]))
            if not name_box:  # Ensure name_box is not None
                return None
            respondent = name_box[5]  
            node_boxes.remove(name_box)  

        node_boxes.sort(key=lambda x: (x[0], x[1], x[2]))  
        return node_boxes, respondent, len(node_boxes)

    def extract_lines(self, slide):
        """
        Extracts all lines (connectors) from a PPT slide.

        :param slide: The slide object.
        :return: Tuple (Dictionary of lines, Total line count)
        """
        lines = {}  

        for i, shape in enumerate(slide.shapes):
            # Check if the shape is a line (includes arrows and connectors)
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                # Explicitly cast shape to Connector type for clarity

                shape: Connector = shape
                x1, y1 = shape.begin_x, shape.begin_y
                x2, y2 = shape.end_x, shape.end_y
                lines[i] = (x1, y1, x2, y2)

        return lines, len(lines)
    
    def extract_slide_image(self, slide):
        """
        Extracts an image from a PowerPoint slide and converts it to an OpenCV (cv2) image.

        This function supports:
        1. Directly inserted images (PICTURE type).
        2. Images inside PLACEHOLDER shapes.
        3. Images inside GROUP shapes.

        :param slide: A PowerPoint slide object.
        :return: An OpenCV image (numpy array) if an image is found, otherwise None.
        """
        for shape in slide.shapes:
            # Check if the shape is a direct image (PICTURE type)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Read the image blob and convert it to a PIL image
                    image_stream = io.BytesIO(shape.image.blob)
                    pil_image = Image.open(image_stream).convert("RGB")

                    # Convert the PIL image to an OpenCV format (numpy array in BGR)
                    return cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)

                except Exception as e:
                    print(f"Error processing image: {e}")
                    return None

            # Check if the shape is a placeholder that contains an image
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "image"):  # Check if placeholder contains an image
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        pil_image = Image.open(image_stream).convert("RGB")
                        return cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)

                    except Exception as e:
                        print(f"Error processing placeholder image: {e}")
                        return None
                    
            # Check if the shape is a group and contains images inside
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Extract and convert the image
                            image_stream = io.BytesIO(sub_shape.image.blob)
                            pil_image = Image.open(image_stream).convert("RGB")
                            return cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)

                        except Exception as e:
                            print(f"Error processing grouped image: {e}")
                            return None

        return None  # No image found in the slide

    def match_lines_to_nodes(self, slide, slide_num):
        """
        Matches each line’s endpoints to the nearest text box.

        :param slide: The slide object.
        :param slide_num: The slide number.
        :return: List of matched relationships (from_text → to_text)
        """
        edges = []  # Stores matched edges (connections) between text boxes
        node_boxes, respondent, node_count = self.extract_node_boxes(slide, slide_num)  # Extract nodes with bounding box coordinates
        lines, line_count = self.extract_lines(slide)  # Extract lines from the slide and count them

        # Create mappings from node_boxes
        node_map = {}  # Maps text → (top, left, right, bottom)
        node_order = {}  # Maps text → index (ranking order)

        for idx, (_, top, left, right, bottom, text) in enumerate(node_boxes):
            node_map[text] = (top, left, right, bottom)  # Store bounding box coordinates for each text label
            node_order[text] = idx  # Assign index based on sorted order

        # Iterate through all detected lines
        for line_id, (x1, y1, x2, y2) in lines.items():
            from_text, to_text = None, None  # Initialize the closest text boxes for each line
            min_distance_from = float("inf")  # Track the minimum distance for the line's start point
            min_distance_to = float("inf")  # Track the minimum distance for the line's end point

            # Iterate over all nodes to find the closest match for both start and end points
            for text, (top, left, right, bottom) in node_map.items():
                # Compute the center point of the text box
                center_x = (left + right) / 2
                center_y = (top + bottom) / 2

                # Calculate Euclidean distances between the line endpoints and the text box center
                dist_from = ((center_x - x1) ** 2 + (center_y - y1) ** 2) ** 0.5
                dist_to = ((center_x - x2) ** 2 + (center_y - y2) ** 2) ** 0.5

                # Update the closest matching text box for the start point
                if dist_from < min_distance_from:
                    from_text = text
                    min_distance_from = dist_from

                # Update the closest matching text box for the end point
                if dist_to < min_distance_to:
                    to_text = text
                    min_distance_to = dist_to

            # Ensure a valid connection (avoid self-connections)
            # Reorder nodes to ensure `from_text` appears lower (has a larger Y-coordinate)
            # PowerPoint uses an inverted Y-axis, meaning larger Y values are lower on the slide
            if from_text and to_text and from_text != to_text:
                from_top, _, _, from_bottom = node_map[from_text]
                to_top, _, _, to_bottom = node_map[to_text]

                # Swap if `to_text` is positioned higher than `from_text`
                if to_top > from_top:
                    from_text, to_text = to_text, from_text
                
                edges.append((from_text, to_text))

        sorted_edges = sorted(edges, key=lambda edge: (node_order.get(edge[0], float("inf")), 
                                                       node_order.get(edge[1], float("inf"))))
        return sorted_edges, respondent, node_count, line_count

    def process_slide(self, slide, slide_num):
        """
        Determines whether to use PPT-based or Image-based extraction.

        :param slide: The slide object.
        :param slide_num: The slide number.
        :return: List of extracted relationships.
        """
        node_boxes, respondent, node_count = self.extract_node_boxes(slide, slide_num)  # Extract nodes with bounding box coordinates
        lines, line_count = self.extract_lines(slide)  # Extract lines from the slide and count them

        if node_count > 0 and line_count > 0:
            return self.match_lines_to_nodes(slide, slide_num)
        
        elif node_count > 0 and line_count == 0:
            return [], respondent, node_count, line_count  # No image found, return empty relationships
        
        elif node_count == 0 and line_count == 0:
            image = self.extract_slide_image(slide)

            if image is not None:
                image_processor = ImageProcessor(image)
                edges = image_processor.match_relationships()
                return edges, respondent, node_count, line_count
            
            else:
                return [], [], node_count, line_count  # No image found, return empty relationships

    def process_pptx_to_dataframe(self):
        """
        Processes all slides and generates a DataFrame.

        :return: DataFrame containing extracted relationships.
        """
        data = []

        for slide_num, slide in enumerate(self.prs.slides, start=1):
            sorted_edges, respondent, node_count, line_count = self.process_slide(slide, slide_num)
            void_edge = int(line_count == 0)

            if not respondent:
                continue
            
            if sorted_edges:
                for from_text, to_text in sorted_edges:
                    data.append([respondent, from_text, to_text, void_edge, slide_num])  # Append relationships to data list   
            else:
                data.append([respondent, None, None, void_edge, slide_num])            

        df = pd.DataFrame(data, columns=['Person Name', 'From', 'To', 'NoMeaningfulEdges', 'Slide Number'])

        id_df_map = pd.read_excel(self.id_path).set_index('full name')['id_ego'].to_dict()
        df.insert(0, 'id_ego', df['Person Name'].map(id_df_map))

        return self.clean_dataframe(df)

    @staticmethod
    def clean_dataframe(df):
        """
        Cleans text data.

        :param df: The input DataFrame.
        :return: Cleaned DataFrame.
        """
        df = df.copy()
        for col in df.select_dtypes(include=['object']).columns:
            df[col] = df[col].astype(str).apply(lambda x: re.sub(r'[\x00-\x1F\x7F]', ' ', x) if pd.notna(x) else x)
        return df